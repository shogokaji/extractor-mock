import pptx

def extract_ppt(filepath: str):
    prs = pptx.Presentation(filepath)
    text = ""

    # 各スライドごとに処理を行う
    for slide in prs.slides:
        # 表はテキストエリアの有無にかかわらず存在するため先に抽出処理
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                row_count = len(table.rows)
                col_count = len(table.columns)
                for r in range(0, row_count):
                    for c in range(0, col_count):
                        c_text = table.cell(r, c).text.replace('\n', '')
                        if c_text == '':
                            continue
                        text += c_text

            # テキストエリアごとの抽出処理
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    c_text = run.text.replace('\n', '')
                    if c_text == '':
                        continue
                    text += run.text
    return text

print(extract_ppt("./docs/sample.pptx"))